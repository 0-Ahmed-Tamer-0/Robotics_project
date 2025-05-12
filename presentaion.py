from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    # Create presentation object
    prs = Presentation()
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Swarm Intelligence Maze Solver"
    subtitle.text = "ACO vs ABC vs PSO Comparison\nRobotics Path Planning Project"

    # Slide 2: Project Objectives
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Project Objectives"
    content = [
        "• Compare swarm intelligence algorithms for robotic path planning",
        "• Evaluate performance metrics:",
        "  - Path optimality\n  - Convergence speed\n  - Scalability",
        "• Develop visualization system for algorithm behavior",
        "• Optimize hyperparameters for maze navigation",
        "• Provide benchmark results for real-world applications"
    ]
    text_box = slide.placeholders[1]
    text_box.text = "\n".join(content)

    # Slide 3: Algorithms Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Algorithms Overview"
    cols = ["Algorithm", "Inspiration", "Strengths"]
    rows = [
        ["ACO", "Ant Colony Pheromones", "Path optimization, Dynamic environments"],
        ["PSO", "Bird Flocking", "Fast convergence, Simple implementation"],
        ["ABC", "Honey Bee Foraging", "Exploration/Exploitation balance"]
    ]
    
    table = slide.shapes.add_table(len(rows)+1, len(cols), 
        Inches(1), Inches(1.5), Inches(8), Inches(4)).table
    
    # Format table
    for i, col in enumerate(cols):
        table.cell(0, i).text = col
        table.cell(0, i).fill.solid()
        table.cell(0, i).fill.fore_color.rgb = RGBColor(59, 89, 152)
        
    for row_idx, row_data in enumerate(rows, 1):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_data
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(242, 242, 242)

    # Slide 4: Key Features
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Key Features"
    content = [
        "✓ Real-time algorithm visualization",
        "✓ Dynamic parameter tuning system",
        "✓ Comprehensive benchmarking framework",
        "✓ Statistical significance testing",
        "✓ Path optimization techniques",
        "✓ Cross-algorithm comparison metrics"
    ]
    text_box = slide.placeholders[1]
    text_box.text = "\n".join(content)

    # Slide 5: Performance Metrics
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Performance Comparison"
    content = [
        "Metric           ACO       PSO       ABC",
        "Success Rate     100%      100%      100%",
        "Avg Path Length  56.7 ± 7.4      116.0 ± 8.6      174.7 ± 41.6",
        "Avg Time (s)     31.55s ± 37.38       8.19s ± 3.00       1.07s ± 1.04",
        
    ]
    text_box = slide.placeholders[1]
    text_box.text = "\n".join(content)

    # Slide 6: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Conclusion & Applications"
    content = [
        "Key Findings:",
        "- PSO: Fastest convergence for simple mazes",
        "- ACO: Best path optimization in complex environments",
        "- ABC: Balanced exploration for dynamic scenarios",
        "",
        "Applications:",
        "• Autonomous drone navigation",
        "• Warehouse robotics path planning",
        "• Emergency evacuation systems",
        "• Game AI pathfinding"
    ]
    text_box = slide.placeholders[1]
    text_box.text = "\n".join(content)

    # Save presentation
    prs.save('SwarmMazeSolver.pptx')

if __name__ == "__main__":
    create_presentation()